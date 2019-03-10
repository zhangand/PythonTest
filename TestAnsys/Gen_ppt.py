'''
Created on Mar 2, 2019

@author: zhangand
'''
from pptx import Presentation
from pptx import chart
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches, Pt  # Cm
import os
import win32com.client

import string

z_info = []

class gen_ppt:
    def read_csv(self, file_path, filename):
        global z_info
        file_obj = open(str(file_path) + '/' + str(filename), 'r')
        i = 0
        for line in file_obj:
            z_info.append(line)
            z_info[i] = z_info[i].rstrip()
            z_info[i] = z_info[i].split(',')
            i = i + 1
        file_obj.close()

    def do_ppt(self,):
        # 创建幻灯片 ------
        prs = Presentation('template.pptx')

        subject_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(subject_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[10]
        title.text = 'Simulation Report'
        subtitle.text = "Tina Chen"

        title_only_slide_layout = prs.slide_layouts[2]
        slide = prs.slides.add_slide(title_only_slide_layout)
        shapes = slide.shapes

        shapes.title.text = 'Simulation Report'

        currentPath = os.getcwd()
        self.read_csv(currentPath, 'VDD_GPU_Z11_table.csv')

        # 表格样式 --------------------
        rows = len(z_info)
        cols = len(z_info[0])
        width = Inches(0.8)*len(z_info[0])  # Inches(6.0)
        height = Inches(0.3)*len(z_info)  # Inches(0.8)
        top = Inches(1)
        left = Inches(9.3)-width  # Inches(2.0)

        # 添加表格到幻灯片 --------------------
        table = shapes.add_table(rows, cols, left, top, width, height).table

        # 设置单元格宽度
        table.columns[0].width = Inches(1.5)

        # 定义表格数据 ------
        for i in range(0, len(z_info), 1):
            for j in range(0, len(z_info[0]), 1):
                if z_info[i][j].startswith('"'):
                    z_info[i][j] = z_info[i][j].replace('"', "")
                    table.cell(i,j).text = z_info[i][j]
                    cell = table.rows[i].cells[j]
                    paragraph = cell.text_frame.paragraphs[0]
                    paragraph.font.size = Pt(12)
                else:
                    z_info[i][j] = str(round(float(z_info[i][j]), 2))
                    table.cell(i,j).text = z_info[i][j]
                    cell = table.rows[i].cells[j]
                    paragraph = cell.text_frame.paragraphs[0]
                    paragraph.font.size = Pt(12)

        #slide = prs.slides.add_slide(title_only_slide_layout)
        #title = slide.shapes.title
        #title.text = 'Simulation Report'
        top = Inches(4)
        left = Inches(0)
        width = Inches(10)

        pic = slide.shapes.add_picture('./VDD_GPU.jpg', left, top, width=width)

        prs.save('test_template.pptx')

        #ppt = win32com.client.Dispatch('PowerPoint.Application')
        #pptSel = ppt.Presentations.Open(currentPath + '/test_template.pptx', ReadOnly=1, Untitled=0, WithWindow=1)
        #win32com.client.gencache.EnsureDispatch('PowerPoint.Application')
        #slide_count = pptSel.Slides.Count
        #print('slide page number is %d' % (slide_count))
        #ppt.Quit()

    # os.system('POWERPNT.exe')
